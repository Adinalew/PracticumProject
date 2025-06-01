import openai
from gtts import gTTS
from io import BytesIO
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import io
import os
import httpx
from openai import OpenAI, OpenAIError
from dotenv import load_dotenv
from django.conf import settings
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from .models import NoteAudio
from django.core.files.base import ContentFile
from django.conf import settings
import uuid
from pdf2image import convert_from_bytes
from transformers import TrOCRProcessor, VisionEncoderDecoderModel
import torch
from PIL import Image


# load the processor and model just once (at module level)
processor = TrOCRProcessor.from_pretrained("microsoft/trocr-base-handwritten")
model = VisionEncoderDecoderModel.from_pretrained("microsoft/trocr-base-handwritten")

# Hardcoded path to poppler bin
os.environ["PATH"] += os.pathsep + r"C:\poppler\poppler-24.08.0\Library\bin"


# Load environment variables
load_dotenv()

def extract_handwriting_from_image_with_trocr(file_field):
    try:
        print("🔠 Running TrOCR handwriting OCR...")
        # read the image from the uploaded file
        image_bytes = file_field.read()
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")

        # preprocess and predict
        pixel_values = processor(images=image, return_tensors="pt").pixel_values
        generated_ids = model.generate(pixel_values)
        text = processor.batch_decode(generated_ids, skip_special_tokens=True)[0]

        return text.strip()

    except Exception as e:
        print(f"❌ Error in TrOCR: {e}")
        return ""

# Tesseract setup
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def preprocess_image(image_file):
    image = Image.open(image_file)

    # convert to grayscale
    image = image.convert('L')

    # increase contrast more aggressively
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(4.0)  # was 2

    # apply sharpen filter
    image = image.filter(ImageFilter.SHARPEN)

    # resize image to improve small handwriting recognition
    base_width = 1000
    w_percent = (base_width / float(image.size[0]))
    h_size = int((float(image.size[1]) * float(w_percent)))
    image = image.resize((base_width, h_size), Image.LANCZOS)

    return image

def extract_text_from_image(file_field):
    return extract_handwriting_from_image_with_trocr(file_field)

def extract_text_from_pdf(file):
    file.seek(0)
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in doc:
            page_text = page.get_text()
            text += page_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        return ""

def extract_text_from_docx(file):
    file.seek(0)
    doc = Document(file)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs).strip()

def extract_text_from_pptx(file):
    file.seek(0)
    presentation = Presentation(file)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text.append(paragraph.text)
    return "\n".join(full_text).strip()

def extract_text_from_file(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
        return extract_text_from_image(file)
    elif ext == '.pdf':
        return extract_text_from_pdf(file)
    elif ext == '.docx':
        return extract_text_from_docx(file)
    elif ext == '.pptx':
        return extract_text_from_pptx(file)
    else:
        return ""


def extract_handwriting_from_image(image_file):
    try:
        print("🔍 Starting OCR for image...")
        image = preprocess_image(image_file)
        image_bytes = io.BytesIO()
        image.save(image_bytes, format='JPEG')
        print("🧪 Image preprocessed. Running Tesseract...")

        result = pytesseract.image_to_string(image)
        print("📄 OCR result:\n", result)
        return result.strip()
    except Exception as e:
        print(f"❌ OCR failed: {e}")
        return ""


def extract_text_from_uploaded_file(uploaded_file):
    file_field = uploaded_file.file
    ext = os.path.splitext(file_field.name)[1].lower()
    print(f"📂 Received file: {file_field.name} (ext={ext})")

    try:
        if ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            print("🧠 Routing to extract_text_from_image()")
            return extract_text_from_image(file_field)

        elif ext == '.txt':
            print("📄 Reading .txt file...")
            file_field.open()
            content_bytes = file_field.read()
            return content_bytes.decode('utf-8', errors='ignore').strip()

        elif ext == '.pdf':
            print("📄 Routing to extract_text_from_pdf()")
            return extract_text_from_pdf(file_field)

        elif ext == '.docx':
            print("📄 Routing to extract_text_from_docx()")
            return extract_text_from_docx(file_field)

        elif ext == '.pptx':
            print("📄 Routing to extract_text_from_pptx()")
            return extract_text_from_pptx(file_field)

        else:
            print("⚠️ Unknown file type")
            return ""

    except Exception as e:
        print(f"❌ Error in extract_text_from_uploaded_file: {e}")
        return ""

def generate_tts_audio(text):
    tts = gTTS(text)
    audio_stream = BytesIO()
    tts.write_to_fp(audio_stream)
    audio_stream.seek(0)
    return audio_stream.read()

def get_text_from_session(session):
    notes = session.extracted_notes.all()
    return "\n\n".join(note.text for note in notes if note.text.strip())

def generate_study_review(text, options=None):
    api_key = os.getenv("OPENAI_API_KEY")
    cert_path = "C:/certificate/2023 techloq bundle certificate.crt"

    if os.path.exists(cert_path):
        print("✅ Using Techloq certificate for SSL.")
        http_client = httpx.Client(verify=cert_path)
    else:
        print("ℹ️ No custom cert found. Using default SSL.")
        http_client = httpx.Client()

    client = OpenAI(api_key=api_key, http_client=http_client)

    # Include user preferences in prompt
    if options:
        notes = (
            f"Language level: {options.get('language_level')}, "
            f"Teaching style: {options.get('explanation_style')}, "
            f"Depth: {options.get('review_depth')}, "
            f"User notes: {options.get('additional_notes') or 'None'}."
        )
    else:
        notes = "No additional preferences provided."

    prompt = (
        f"You are an AI assistant. Based on the following notes, create a study review.\n"
        f"{notes}\n\n"
        f"Study text:\n{text}"
    )

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": prompt}
            ]
        )
        return response.choices[0].message.content.strip()
    except OpenAIError as e:
        print(f"❌ OpenAI API error: {e}")
        return "An error occurred while generating the study review."

def generate_followup_response(question, previous_review, options=None):
    api_key = os.getenv("OPENAI_API_KEY")
    cert_path = "C:/certificate/2023 techloq bundle certificate.crt"

    if os.path.exists(cert_path):
        http_client = httpx.Client(verify=cert_path)
    else:
        http_client = httpx.Client()

    client = OpenAI(api_key=api_key, http_client=http_client)

    style = options.get('explanation_style', 'normal') if options else 'normal'

    prompt = (
        f"You are a helpful tutor. Based on the study material below, respond to the user's follow-up question.\n"
        f"Make sure to answer in the appropriate tone. Explanation style: {style}.\n\n"
        f"Study Material:\n{previous_review}\n\n"
        f"Follow-up Question:\n{question}"
    )

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a kind, helpful tutor."},
                {"role": "user", "content": prompt}
            ]
        )
        return response.choices[0].message.content.strip()
    except OpenAIError as e:
        print(f"Error generating follow-up: {e}")
        return "Sorry, something went wrong while generating your answer."

def generate_flashcards_from_text(text, custom_prompt=None):
    api_key = os.getenv("OPENAI_API_KEY")
    cert_path = "C:/certificate/2023 techloq bundle certificate.crt"

    if not custom_prompt:
        custom_prompt = (
            "Read the following study notes and generate helpful flashcards. "
            "They should include important terms and clear explanations. "
            "Return the result in this format:\n\n"
            "Q: ...\nA: ...\n\nQ: ...\nA: ..."
        )

    # Always append strict instruction at the end
    enforced_format = (
    "\n\n---\nIMPORTANT: Return only flashcards in this format:\n"
    "Q: ...\nA: ...\n\nQ: ...\nA: ...\n"
    "Do NOT include explanations, greetings, or any markdown formatting."
    )

    prompt = f"{custom_prompt}\n\nSTUDY TEXT:\n{text}{enforced_format}"

    try:
        http_client = httpx.Client(verify=cert_path) if os.path.exists(cert_path) else httpx.Client()
        client = OpenAI(api_key=api_key, http_client=http_client)

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful tutor."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
        )

        content = response.choices[0].message.content
        print("🔎 OpenAI returned:\n", content)  # 👈 TEMP DEBUGGING LINE

        # Parse into Q/A
        cards = []
        lines = content.split('\n')
        question = answer = None
        for line in lines:
            if line.strip().startswith("Q:"):
                question = line.strip()[2:].strip()
            elif line.strip().startswith("A:"):
                answer = line.strip()[2:].strip()
                if question:
                    cards.append((question, answer))
                    question = answer = None

        print(f"✅ Parsed {len(cards)} flashcards.")  # 👈 TEMP DEBUGGING LINE
        return cards

    except OpenAIError as e:
        print(f"❌ OpenAI API error: {e}")
        return []
    except Exception as e:
        print(f"❌ General error generating flashcards: {e}")
        return []

def generate_note_audio(note):
    """
    Uses OpenAI TTS to generate natural-sounding voice for a single ExtractedNote.
    Saves the audio file and returns the NoteAudio instance.
    """
    api_key = os.getenv("OPENAI_API_KEY")
    cert_path = "C:/certificate/2023 techloq bundle certificate.crt"

    if os.path.exists(cert_path):
        http_client = httpx.Client(verify=cert_path)
    else:
        http_client = httpx.Client()

    client = OpenAI(api_key=api_key, http_client=http_client)

    try:
        response = client.audio.speech.create(
            model="tts-1",  # Or use "tts-1-hd" if allowed
            voice="shimmer",  # options: alloy, echo, fable, onyx, nova, shimmer
            input=note.text[:4096]  # OpenAI TTS limit per request
        )

        mp3_data = response.content
        filename = f"note_{note.id}_{uuid.uuid4().hex[:8]}.mp3"

        note_audio = NoteAudio(note=note)
        note_audio.audio_file.save(filename, ContentFile(mp3_data))
        note_audio.save()

        return note_audio

    except Exception as e:
        print(f"❌ Failed to generate AI voice for note {note.id}: {e}")
        return None

def extract_handwriting_from_pdf(file):
    try:
        file.seek(0)
        images = convert_from_bytes(file.read())
        text = ""
        for img in images:
            img = preprocess_image(img)
            text += pytesseract.image_to_string(img) + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error in handwriting PDF OCR: {e}")
        return ""